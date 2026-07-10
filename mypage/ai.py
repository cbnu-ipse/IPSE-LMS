import json

from django.conf import settings

def extract_text(file_obj, filename):
    """PDF/DOCX/PPTX 파일 객체에서 본문 텍스트를 뽑는다. 실패하면 빈 문자열을 반환한다."""
    name = filename.lower()
    try:
        if name.endswith(".pdf"):
            return _extract_pdf(file_obj)
        if name.endswith(".docx") or name.endswith(".doc"):
            return _extract_docx(file_obj)
        if name.endswith(".pptx") or name.endswith(".ppt"):
            return _extract_pptx(file_obj)
    except Exception:
        return ""
    return ""


def _extract_pdf(f):
    from pypdf import PdfReader

    reader = PdfReader(f)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(f):
    import docx

    document = docx.Document(f)
    return "\n".join(p.text for p in document.paragraphs)


def _extract_pptx(f):
    from pptx import Presentation

    prs = Presentation(f)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)


def _client():
    from openai import OpenAI

    if not settings.OPENAI_API_KEY:
        raise RuntimeError("OPENAI_API_KEY가 설정되어 있지 않습니다. 운영 .env에 키를 추가한 뒤 다시 시도해주세요.")
    return OpenAI(api_key=settings.OPENAI_API_KEY)


def generate_summary(extracted_text):
    """업로드된 문서 본문을 길고 구조화된 요약으로 반환한다. 본문이 없으면 호출 없이 빈 문자열."""
    if not extracted_text.strip():
        return ""

    client = _client()
    prompt = (
        "다음은 학생이 올린 학습 자료 본문입니다. 시험 대비용으로 충분히 활용할 수 있도록 길고 "
        "자세하게 요약해주세요.\n"
        "- 내용을 몇 개의 소제목(##)으로 나누고, 각 소제목 아래에 핵심 내용을 문장으로 정리하세요.\n"
        "- 전체 분량은 최소 10문장 이상으로, 중요한 개념/정의/예시를 빠짐없이 포함하세요.\n"
        "- 설명 없이 요약문만 출력하세요.\n\n"
        "--- 학습 자료 본문 ---\n"
        f"{extracted_text[:8000]}"
    )
    response = client.chat.completions.create(
        model=settings.MYPAGE_CHAT_MODEL,
        messages=[{"role": "user", "content": prompt}],
    )
    return response.choices[0].message.content.strip()


QUESTION_TYPE_LABELS = {"ox": "OX(참/거짓)", "short": "단답형", "essay": "서술형"}


def generate_one_question(extracted_text, question_type, existing_questions=()):
    """업로드된 문서 본문을 근거로 지정된 유형(ox/short/essay)의 문제를 1개 생성해 dict로 반환한다.
    existing_questions에 담긴 문제와는 겹치지 않게 요청한다."""
    client = _client()
    avoid_section = ""
    if existing_questions:
        listed = "\n".join(f"- {q}" for q in existing_questions)
        avoid_section = f"\n다음 문제들과는 겹치지 않는 새로운 문제로 만들어주세요:\n{listed}\n"

    label = QUESTION_TYPE_LABELS[question_type]
    prompt = (
        "다음은 학생이 올린 학습 자료 본문입니다. 이 내용만 근거로 복습 문제를 만들어주세요.\n"
        f"- {label} 문제 1개\n"
        f"{avoid_section}\n"
        '아래 JSON 형식으로만 응답하세요: {"question": "...", "answer": "..."}\n\n'
        "--- 학습 자료 본문 ---\n"
        f"{extracted_text[:8000]}"
    )

    response = client.chat.completions.create(
        model=settings.MYPAGE_CHAT_MODEL,
        messages=[{"role": "user", "content": prompt}],
        response_format={"type": "json_object"},
    )
    data = json.loads(response.choices[0].message.content)
    return {"question": data.get("question", ""), "answer": data.get("answer", "")}


def explain_ox_answer(question_text, correct_answer):
    """OX 문제를 틀렸을 때, 정답 근거를 한두 문장으로 설명해준다."""
    client = _client()
    prompt = (
        "다음 OX 문제의 정답이 왜 그런지 학생에게 한두 문장으로 간결하게 설명해주세요.\n\n"
        f"문제: {question_text}\n"
        f"정답: {correct_answer}\n\n"
        "설명만 출력하세요."
    )
    response = client.chat.completions.create(
        model=settings.MYPAGE_CHAT_MODEL,
        messages=[{"role": "user", "content": prompt}],
    )
    return response.choices[0].message.content.strip()


def grade_answer(question_text, correct_answer, user_answer):
    """단답형/서술형 답안을 채점한다. OX는 문자열 비교로 충분해 이 함수를 거치지 않는다."""
    client = _client()
    prompt = (
        "다음 문제에 대한 학생의 답안을 채점해주세요. 완전히 같은 문장이 아니어도 핵심 내용이 "
        "모범답안과 일치하면 정답으로 관대하게 처리하세요.\n\n"
        f"문제: {question_text}\n"
        f"모범답안: {correct_answer}\n"
        f"학생 답안: {user_answer}\n\n"
        '아래 JSON 형식으로만 응답하세요: {"is_correct": true/false, "feedback": "한두 문장 피드백"}'
    )
    response = client.chat.completions.create(
        model=settings.MYPAGE_CHAT_MODEL,
        messages=[{"role": "user", "content": prompt}],
        response_format={"type": "json_object"},
    )
    data = json.loads(response.choices[0].message.content)
    return {"is_correct": bool(data.get("is_correct")), "feedback": data.get("feedback", "")}
